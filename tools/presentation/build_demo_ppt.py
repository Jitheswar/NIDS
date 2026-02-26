#!/usr/bin/env python3
"""Build a cyber-modern PowerPoint deck from Demo.md."""

from __future__ import annotations

import argparse
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


def rgb(hex_color: str) -> RGBColor:
    """Convert #RRGGBB to RGBColor."""
    value = hex_color.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


PALETTE = {
    "bg": rgb("#0A1020"),
    "surface": rgb("#121A2E"),
    "accent_primary": rgb("#00A3FF"),
    "accent_secondary": rgb("#00E0C7"),
    "warning": rgb("#FFB547"),
    "text_primary": rgb("#F5F8FF"),
    "text_secondary": rgb("#B8C2D9"),
}

FONT_TITLE = "Noto Sans"
FONT_BODY = "Noto Sans"
FONT_CODE = "Liberation Mono"


class DeckBuilder:
    """Helper for creating consistent slides."""

    def __init__(self, source_md: Path, output_pptx: Path) -> None:
        self.source_md = source_md
        self.output_pptx = output_pptx
        self.source_text = source_md.read_text(encoding="utf-8")
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(13.333)
        self.presentation.slide_height = Inches(7.5)
        self.live_script = self._extract_live_demo_script()

    def _extract_live_demo_script(self) -> str:
        match = re.search(
            r"## 6\).*?```bash\n(.*?)```",
            self.source_text,
            flags=re.DOTALL,
        )
        if not match:
            raise ValueError("Could not find live demo script block in Demo.md")
        script = match.group(1).strip("\n")
        return script

    def _new_slide(self):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = PALETTE["bg"]
        bg.line.fill.background()

        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = PALETTE["accent_primary"]
        top_bar.line.fill.background()
        return slide

    def _add_title(self, slide, title: str, subtitle: str | None = None) -> None:
        title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.42), Inches(12.0), Inches(0.9))
        title_tf = title_box.text_frame
        title_tf.clear()
        title_tf.word_wrap = True
        p = title_tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_TITLE
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = PALETTE["text_primary"]

        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.2), Inches(12.0), Inches(0.6))
            sub_tf = sub_box.text_frame
            sub_tf.clear()
            sub_tf.word_wrap = True
            sp = sub_tf.paragraphs[0]
            sp.text = subtitle
            sp.font.name = FONT_BODY
            sp.font.size = Pt(18)
            sp.font.color.rgb = PALETTE["text_secondary"]

    def _add_card(self, slide, x: float, y: float, w: float, h: float, title: str | None = None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = PALETTE["surface"]
        card.line.color.rgb = PALETTE["accent_primary"]
        card.line.width = Pt(1.0)

        if title:
            header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.35))
            header.fill.solid()
            header.fill.fore_color.rgb = PALETTE["accent_primary"]
            header.line.fill.background()

            tf = header.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = FONT_BODY
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = PALETTE["bg"]
            p.alignment = PP_ALIGN.LEFT
        return card

    def _add_bullets(
        self,
        slide,
        bullets: list[str],
        x: float,
        y: float,
        w: float,
        h: float,
        size: int = 20,
        color: RGBColor = PALETTE["text_primary"],
        level: int = 0,
    ) -> None:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(4)
        tf.margin_bottom = Pt(4)
        for idx, line in enumerate(bullets):
            para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            para.text = line
            para.level = level
            para.font.name = FONT_BODY
            para.font.size = Pt(size)
            para.font.color.rgb = color

    def _add_code_panel(self, slide, text: str, x: float, y: float, w: float, h: float, font_size: float = 11.0):
        panel = self._add_card(slide, x, y, w, h)
        panel.line.color.rgb = PALETTE["accent_secondary"]
        code_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.14), Inches(w - 0.36), Inches(h - 0.28))
        tf = code_box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = FONT_CODE
        p.font.size = Pt(font_size)
        p.font.color.rgb = PALETTE["text_primary"]

    def _add_note(self, slide, lines: list[str]) -> None:
        notes = slide.notes_slide.notes_text_frame
        notes.clear()
        for idx, line in enumerate(lines):
            p = notes.paragraphs[0] if idx == 0 else notes.add_paragraph()
            p.text = line

    def slide_1_title(self) -> None:
        slide = self._new_slide()
        self._add_title(slide, "AI-Enabled NIDS Backend Demo", "Cyber-Modern Classroom Evaluation Deck")

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.9), Inches(4.8), Inches(0.06))
        accent.fill.solid()
        accent.fill.fore_color.rgb = PALETTE["accent_secondary"]
        accent.line.fill.background()

        self._add_card(slide, 0.65, 2.2, 5.8, 3.8, "Presenter Details")
        self._add_bullets(
            slide,
            [
                "Name: __________________________",
                "Roll No: ________________________",
                "Class/Section: ___________________",
                "Date: __________________________",
            ],
            0.85,
            2.75,
            5.4,
            2.9,
            size=20,
        )

        self._add_card(slide, 6.7, 2.2, 5.95, 3.8, "Demo Snapshot")
        self._add_bullets(
            slide,
            [
                "Duration: 10 to 15 minutes",
                "Mode: Live API + terminal workflow",
                "Core focus: Auth, sensors, audit, AI analytics",
                "Target: Clear, defensible security architecture",
            ],
            6.95,
            2.75,
            5.5,
            2.9,
            size=18,
        )

        self._add_note(
            slide,
            [
                "Introduce the project as an AI-enabled NIDS backend demo.",
                "Set expectation: live, end-to-end verification in 10 to 15 minutes.",
                "Mention that the focus is security controls and observable outcomes.",
            ],
        )

    def slide_2_goals(self) -> None:
        slide = self._new_slide()
        self._add_title(slide, "Demo Goals and Outcomes", "What the evaluator should see clearly")

        self._add_card(slide, 0.65, 1.9, 6.2, 4.9, "Goals")
        self._add_bullets(
            slide,
            [
                "Demonstrate reliable startup and readiness checks",
                "Prove authentication and RBAC with Keycloak-backed JWT",
                "Walk through sensor onboarding lifecycle",
                "Show tamper-aware audit visibility",
                "Show operational security checks and anomaly interfaces",
            ],
            0.9,
            2.35,
            5.8,
            4.2,
            size=18,
        )

        self._add_card(slide, 7.05, 1.9, 5.6, 4.9, "Evaluator Outcomes")
        self._add_bullets(
            slide,
            [
                "System is healthy and production-shaped",
                "Access control is role enforced, not just login based",
                "Sensor enrollment uses single-use key flow",
                "Critical actions appear in audit logs",
                "Security monitoring endpoints are present and queryable",
            ],
            7.3,
            2.35,
            5.1,
            4.2,
            size=17,
        )

        self._add_note(
            slide,
            [
                "Frame this slide as the grading checklist.",
                "Tell the evaluator each item will be demonstrated with concrete API evidence.",
                "Signal that outcomes are mapped to live command outputs.",
            ],
        )

    def slide_3_scope(self) -> None:
        slide = self._new_slide()
        self._add_title(slide, "What Is Being Demonstrated", "Core capabilities in this session")

        cards = [
            ("Identity and Access", "Keycloak login/session flow\nJWT claims and role checks"),
            ("Sensor Lifecycle", "Create sensor\nEnroll with single-use key\nActivate"),
            ("Audit Pipeline", "Audit events recorded\nQuery recent logs"),
            ("Security Operations", "Run rotation health checks\nReview findings"),
            ("AI Analytics Surface", "Inspect anomaly and risk-score endpoints"),
        ]

        x_positions = [0.65, 4.55, 8.45, 0.65, 6.55]
        y_positions = [2.0, 2.0, 2.0, 4.4, 4.4]
        widths = [3.7, 3.7, 4.2, 5.5, 6.1]
        heights = [2.15, 2.15, 2.15, 2.25, 2.25]

        for idx, (title, body) in enumerate(cards):
            self._add_card(slide, x_positions[idx], y_positions[idx], widths[idx], heights[idx], title)
            self._add_bullets(
                slide,
                body.split("\n"),
                x_positions[idx] + 0.2,
                y_positions[idx] + 0.55,
                widths[idx] - 0.35,
                heights[idx] - 0.65,
                size=14,
            )

        self._add_note(
            slide,
            [
                "Use this as a scope boundary: these are in-demo capabilities.",
                "State that deeper features exist but are out of current time window.",
                "Transition to architecture next to show how components support the flow.",
            ],
        )

    def slide_4_architecture(self) -> None:
        slide = self._new_slide()
        self._add_title(slide, "Architecture Overview", "FastAPI core with security-first service mesh")

        app_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(2.2), Inches(4.9), Inches(2.8))
        app_box.fill.solid()
        app_box.fill.fore_color.rgb = PALETTE["surface"]
        app_box.line.color.rgb = PALETTE["accent_primary"]
        app_tf = app_box.text_frame
        app_tf.clear()
        app_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = app_tf.paragraphs[0]
        p.text = "FastAPI App (:8000)\nAuth | Sensors | Security | Audit"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_BODY
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = PALETTE["text_primary"]

        nodes = [
            ("Sensors", 0.85, 2.6, 2.5, 1.4, PALETTE["accent_secondary"]),
            ("User/Swagger", 0.85, 4.35, 2.5, 1.4, PALETTE["accent_secondary"]),
            ("Keycloak", 10.0, 1.8, 2.4, 1.1, PALETTE["accent_primary"]),
            ("MariaDB", 10.0, 3.05, 2.4, 1.1, PALETTE["accent_primary"]),
            ("Redis", 10.0, 4.3, 2.4, 1.1, PALETTE["accent_primary"]),
            ("Loki", 10.0, 5.55, 2.4, 1.1, PALETTE["accent_primary"]),
        ]
        for label, x, y, w, h, line_color in nodes:
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            box.fill.solid()
            box.fill.fore_color.rgb = PALETTE["surface"]
            box.line.color.rgb = line_color
            tf = box.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = tf.paragraphs[0]
            para.text = label
            para.alignment = PP_ALIGN.CENTER
            para.font.name = FONT_BODY
            para.font.size = Pt(16)
            para.font.color.rgb = PALETTE["text_primary"]

        line_specs = [
            (3.35, 3.3, 0.85, 3.3),
            (3.35, 5.0, 0.85, 5.0),
            (9.1, 2.35, 10.0, 2.35),
            (9.1, 3.6, 10.0, 3.6),
            (9.1, 4.85, 10.0, 4.85),
            (9.1, 6.1, 10.0, 6.1),
        ]
        for x1, y1, x2, y2 in line_specs:
            conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
            conn.line.color.rgb = PALETTE["accent_secondary"]
            conn.line.width = Pt(1.8)

        footer = slide.shapes.add_textbox(Inches(0.65), Inches(6.72), Inches(12.1), Inches(0.45))
        ft = footer.text_frame
        ft.clear()
        fp = ft.paragraphs[0]
        fp.text = "Private network: MariaDB, Redis, Loki, step-ca, Infisical | Frontend network: App, Keycloak"
        fp.font.name = FONT_BODY
        fp.font.size = Pt(12)
        fp.font.color.rgb = PALETTE["text_secondary"]
        fp.alignment = PP_ALIGN.CENTER

        self._add_note(
            slide,
            [
                "Explain the app as the orchestration point between identity, state, and logs.",
                "Highlight private/internal services versus externally exposed surfaces.",
                "Call out that architecture supports zero-trust and observability.",
            ],
        )

    def slide_5_defense_roles(self) -> None:
        slide = self._new_slide()
        self._add_title(slide, "Defense-in-Depth and RBAC Roles", "Layered controls and bounded permissions")

        self._add_card(slide, 0.65, 1.9, 7.05, 4.9, "Security Layers")
        self._add_bullets(
            slide,
            [
                "Layer 1: Keycloak SSO, MFA-ready auth, role validation",
                "Layer 2: Sensor onboarding via single-use enrollment keys",
                "Layer 3: Audit logging to DB and Loki pipeline",
                "Layer 4: Rotation health checks and anomaly endpoints",
                "Layer 5: Optional mTLS + ZTNA enforcement controls",
            ],
            0.95,
            2.35,
            6.55,
            4.2,
            size=16,
        )

        self._add_card(slide, 7.95, 1.9, 4.7, 4.9, "Role Summary")
        self._add_bullets(
            slide,
            [
                "super_admin",
                "security_analyst",
                "auditor",
                "sensor_manager",
            ],
            8.2,
            2.35,
            4.2,
            1.9,
            size=17,
        )
        self._add_bullets(
            slide,
            [
                "Super admin: full control + rotation checks",
                "Security analyst: anomaly + risk views",
                "Auditor: read audit evidence",
                "Sensor manager: sensor lifecycle operations",
            ],
            8.2,
            4.2,
            4.2,
            2.3,
            size=13,
            color=PALETTE["text_secondary"],
        )

        self._add_note(
            slide,
            [
                "Use this slide to justify principle-of-least-privilege design.",
                "Clarify that role checks are enforced server-side per endpoint.",
                "Tie role outputs to the upcoming /auth/me proof.",
            ],
        )

    def slide_6_setup(self) -> None:
        slide = self._new_slide()
        self._add_title(slide, "Demo Setup Checklist", "Pre-demo reliability steps")

        self._add_card(slide, 0.65, 1.9, 6.3, 4.95, "Environment and Prerequisites")
        self._add_bullets(
            slide,
            [
                "Install: Docker, Docker Compose, curl, jq",
                "Copy env: cp .env.example .env",
                "Required: DB_USER=nids",
                "Required: KEYCLOAK_CLIENT_SECRET=change_me_client_secret",
                "Use clean startup before evaluation day",
            ],
            0.95,
            2.35,
            5.95,
            4.2,
            size=16,
        )

        self._add_card(slide, 7.15, 1.9, 5.5, 4.95, "Clean Startup Commands")
        self._add_code_panel(
            slide,
            "docker compose down -v\n"
            "docker compose up -d --build\n\n"
            "docker compose ps\n"
            "curl -sS http://localhost:8000/health/ready | jq\n\n"
            "Open:\n"
            "http://localhost:8000/docs\n"
            "http://localhost:8080",
            7.35,
            2.35,
            5.1,
            4.25,
            font_size=12,
        )

        self._add_note(
            slide,
            [
                "State that this checklist prevents almost all demo-day failures.",
                "Emphasize two fragile values: DB_USER and Keycloak client secret.",
                "Mention that readiness endpoint is the go/no-go gate before live flow.",
            ],
        )

    def slide_7_flow(self) -> None:
        slide = self._new_slide()
        self._add_title(slide, "Live Demo Flow Timeline", "Nine steps from readiness to monitoring")

        steps = [
            "1. Verify /health/ready",
            "2. POST /auth/login",
            "3. GET /auth/me",
            "4. POST /sensors/",
            "5. POST /sensors/{id}/enroll",
            "6. POST /sensors/activate",
            "7. GET /sensors/",
            "8. GET /audit/logs?limit=5",
            "9. POST /security/rotation-health/run",
        ]

        left = 0.85
        top = 2.0
        chip_w = 4.0
        chip_h = 0.86
        x_gap = 0.4
        y_gap = 0.3

        for idx, step in enumerate(steps):
            row = idx // 3
            col = idx % 3
            x = left + col * (chip_w + x_gap)
            y = top + row * (chip_h + y_gap)
            chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(chip_w), Inches(chip_h))
            chip.fill.solid()
            chip.fill.fore_color.rgb = PALETTE["surface"]
            chip.line.color.rgb = PALETTE["accent_secondary"] if idx % 2 == 0 else PALETTE["accent_primary"]
            tf = chip.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = step
            p.font.name = FONT_BODY
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = PALETTE["text_primary"]
            p.alignment = PP_ALIGN.CENTER

        self._add_card(slide, 0.85, 5.2, 12.05, 1.45, "Narrative Anchor")
        self._add_bullets(
            slide,
            [
                "Authenticate -> Authorize -> Provision -> Activate -> Observe -> Validate Security Health",
            ],
            1.1,
            5.75,
            11.5,
            0.7,
            size=18,
        )

        self._add_note(
            slide,
            [
                "Walk quickly through each numbered step and expected output.",
                "Tell evaluator this order minimizes context switching during terminal demo.",
                "Keep this as the roadmap slide to return to if interrupted.",
            ],
        )

    def slide_8_key_commands(self) -> None:
        slide = self._new_slide()
        self._add_title(slide, "Key Command Highlights", "Concise commands on the main flow")

        self._add_code_panel(
            slide,
            "curl -sS http://localhost:8000/health/ready | jq",
            0.65,
            2.0,
            12.0,
            1.05,
            font_size=14,
        )
        self._add_code_panel(
            slide,
            "TOKEN=$(curl -sS -X POST http://localhost:8000/auth/login \\\n"
            "  -H 'Content-Type: application/json' \\\n"
            "  -d '{\"username\":\"admin\",\"password\":\"Admin@nids2024!\"}' | jq -r '.access_token')",
            0.65,
            3.2,
            12.0,
            1.55,
            font_size=10.5,
        )
        self._add_code_panel(
            slide,
            "curl -sS -X POST http://localhost:8000/sensors/ \\\n"
            "  -H \"Authorization: Bearer $TOKEN\" \\\n"
            "  -H 'Content-Type: application/json' \\\n"
            "  -d '{\"name\":\"demo-sensor\",\"network_segment\":\"172.28.10.0/24\"}'",
            0.65,
            4.95,
            12.0,
            1.55,
            font_size=10.5,
        )

        self._add_note(
            slide,
            [
                "Show only high-signal commands here to keep attention on the logic.",
                "Mention full script is available in appendix slide and Demo.md.",
                "Pause briefly after login to validate token and identity context.",
            ],
        )

    def slide_9_results(self) -> None:
        slide = self._new_slide()
        self._add_title(slide, "Proof of Results", "What successful outputs look like")

        self._add_card(slide, 0.65, 1.9, 6.15, 4.95, "Validation Evidence")
        self._add_bullets(
            slide,
            [
                "GET /auth/me shows role claim: super_admin",
                "Sensor status progression: pending -> active",
                "GET /sensors/ returns updated inventory and totals",
                "GET /audit/logs returns latest event types",
                "POST /security/rotation-health/run returns counts",
            ],
            0.95,
            2.35,
            5.75,
            4.2,
            size=15,
        )

        self._add_card(slide, 6.95, 1.9, 5.7, 4.95, "Expected Signals")
        self._add_code_panel(
            slide,
            "\"status\": \"ready\"\n"
            "\"roles\": [\"super_admin\"]\n"
            "\"latest_event\": \"sensor_activated\"\n"
            "\"critical\": <n>, \"warning\": <n>",
            7.2,
            2.35,
            5.2,
            2.2,
            font_size=13,
        )
        self._add_bullets(
            slide,
            [
                "Evaluator should see evidence, not only claims.",
                "Every key action has matching system output.",
            ],
            7.2,
            4.8,
            5.2,
            1.6,
            size=14,
            color=PALETTE["text_secondary"],
        )

        self._add_note(
            slide,
            [
                "Use this slide to connect commands to measurable proof points.",
                "Point out role claim and sensor state change as core correctness checks.",
                "Reinforce that audit trail confirms accountability.",
            ],
        )

    def slide_10_ai(self) -> None:
        slide = self._new_slide()
        self._add_title(slide, "AI and Monitoring Endpoints", "Security analytics visibility")

        self._add_card(slide, 0.65, 2.0, 12.0, 1.05, "Endpoints to show in Swagger")
        self._add_code_panel(
            slide,
            "GET /security/anomalies\n"
            "GET /security/anomalies/risk-score\n"
            "POST /security/rotation-health/run",
            0.85,
            2.22,
            11.6,
            0.7,
            font_size=15,
        )

        self._add_card(slide, 0.65, 3.35, 5.8, 3.45, "How to explain AI in demo")
        self._add_bullets(
            slide,
            [
                "Anomaly service evaluates event streams",
                "Risk scoring synthesizes weighted signals",
                "Endpoints expose current analytic state",
                "Demonstrates AI-backed security operations",
            ],
            0.9,
            3.8,
            5.35,
            2.75,
            size=15,
        )

        self._add_card(slide, 6.8, 3.35, 5.85, 3.45, "Monitoring posture")
        self._add_bullets(
            slide,
            [
                "Operational health checks for rotations",
                "Findings classified by severity",
                "Early warning before cert/secret failures",
                "Works with audit trail for incident response",
            ],
            7.05,
            3.8,
            5.35,
            2.75,
            size=15,
        )

        self._add_note(
            slide,
            [
                "Answer 'Where is AI?' by pointing to anomaly and risk-score endpoints.",
                "Describe analytics as part of security decision support, not marketing claim.",
                "Mention rotation checks as operational hardening complement.",
            ],
        )

    def slide_11_issues(self) -> None:
        slide = self._new_slide()
        self._add_title(slide, "Common Issues and Quick Fixes", "Demo-day recovery playbook")

        issues = [
            (
                "App not starting",
                "Check .env values (DB_USER=nids and Keycloak secret); run clean restart",
            ),
            (
                "/auth/login returns 423 Locked",
                "Clear lockout keys in Redis for admin before login retry",
            ),
            (
                "Services unhealthy",
                "Use compose status and logs for keycloak/mariadb/redis/loki",
            ),
            (
                "step-ca warning in rotation check",
                "Expected if STEP_CA_FINGERPRINT unset and mTLS enrollment not shown",
            ),
        ]

        y = 2.0
        for idx, (issue, fix) in enumerate(issues):
            card = self._add_card(slide, 0.65, y, 12.0, 1.05)
            card.line.color.rgb = PALETTE["warning"] if idx in (0, 1) else PALETTE["accent_primary"]
            self._add_bullets(slide, [issue], 0.9, y + 0.14, 3.7, 0.34, size=13)
            self._add_bullets(slide, [fix], 4.35, y + 0.14, 8.0, 0.62, size=12, color=PALETTE["text_secondary"])
            y += 1.17

        self._add_code_panel(
            slide,
            "docker compose logs --tail=100 app\n"
            "docker compose ps\n"
            "docker compose logs --tail=150 keycloak mariadb redis loki",
            0.65,
            6.72,
            12.0,
            0.7,
            font_size=9.5,
        )

        self._add_note(
            slide,
            [
                "Tell evaluator you planned failure handling before demo day.",
                "If a failure occurs, use this playbook without losing presentation flow.",
                "Keep remediation short, then return to main timeline slide.",
            ],
        )

    def slide_12_appendix(self) -> None:
        slide = self._new_slide()
        self._add_title(slide, "Appendix: Full Live Demo Script", "Copy/paste fallback from Demo.md")

        script_lines = self.live_script.splitlines()
        split_at = len(script_lines) // 2
        left_text = "\n".join(script_lines[:split_at])
        right_text = "\n".join(script_lines[split_at:])

        self._add_code_panel(slide, left_text, 0.65, 1.95, 6.1, 5.3, font_size=6.5)
        self._add_code_panel(slide, right_text, 6.95, 1.95, 5.7, 5.3, font_size=6.5)

        self._add_note(
            slide,
            [
                "This slide mirrors the complete script from section 6 of Demo.md.",
                "Use this only if you need to recover quickly during live demonstration.",
                "Normal flow should rely on concise command slides plus terminal execution.",
            ],
        )

    def build(self) -> None:
        self.slide_1_title()
        self.slide_2_goals()
        self.slide_3_scope()
        self.slide_4_architecture()
        self.slide_5_defense_roles()
        self.slide_6_setup()
        self.slide_7_flow()
        self.slide_8_key_commands()
        self.slide_9_results()
        self.slide_10_ai()
        self.slide_11_issues()
        self.slide_12_appendix()

        self.output_pptx.parent.mkdir(parents=True, exist_ok=True)
        self.presentation.save(str(self.output_pptx))


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate cyber-modern NIDS demo PPTX.")
    parser.add_argument("--source", type=Path, required=True, help="Path to Demo.md source file.")
    parser.add_argument("--out", type=Path, required=True, help="Output .pptx path.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    builder = DeckBuilder(args.source, args.out)
    builder.build()
    print(f"Generated presentation: {args.out}")


if __name__ == "__main__":
    main()
